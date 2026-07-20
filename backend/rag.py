import os
import logging

# Must be set before any numpy/openblas import
os.environ['OPENBLAS_NUM_THREADS'] = '4'
os.environ['OMP_NUM_THREADS'] = '4'

from langchain_community.document_loaders import Docx2txtLoader, TextLoader
from langchain_core.documents import Document
from pptx import Presentation
import pypdfium2 as pdfium
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate
from dotenv import load_dotenv

load_dotenv()

logger = logging.getLogger("quill_ai.rag")

# ── ChromaDB path ─────────────────────────────────────────────────────────────
CHROMA_PATH = os.environ.get("CHROMA_PATH", "chroma_db")

# ── L2 distance threshold for all-MiniLM-L6-v2 ───────────────────────────────
SIMILARITY_THRESHOLD = float(os.environ.get("SIMILARITY_THRESHOLD", "1.5"))

# ── Embeddings (loaded once at module import) ─────────────────────────────────
logger.info("Loading HuggingFace embeddings model…")
embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
logger.info("Embeddings model loaded.")

# ── ChromaDB singleton ────────────────────────────────────────────────────────
_db: Chroma | None = None


def init_db() -> Chroma:
    global _db
    if _db is None:
        logger.info("Initialising ChromaDB at path: %s", CHROMA_PATH)
        _db = Chroma(persist_directory=CHROMA_PATH, embedding_function=embeddings)
        logger.info("ChromaDB ready.")
    return _db


def get_db() -> Chroma:
    global _db
    if _db is None:
        init_db()
    return _db


# ── Document processing ───────────────────────────────────────────────────────
def process_document(filepath: str, filename: str, session_id: str) -> None:
    """Load a document, split into chunks, and add to ChromaDB — tagged with session_id."""
    logger.info("process_document: loading '%s' from '%s' (session=%s)", filename, filepath, session_id)

    if filename.lower().endswith(".pdf"):
        # Extract text page-by-page ourselves instead of relying on
        # PyPDFium2Loader. That LangChain wrapper can silently switch between
        # "one Document per page" and "one Document for the whole file"
        # depending on the installed version, and in the latter mode there's
        # no per-page metadata at all — every chunk then falls back to
        # page=0, which is exactly the "Page 0 for everything" bug this
        # replaces. Doing it directly guarantees correct 1-indexed pages.
        pdf = pdfium.PdfDocument(filepath)
        docs = []
        try:
            for i in range(len(pdf)):
                page = pdf[i]
                textpage = page.get_textpage()
                try:
                    text = textpage.get_text_range().strip()
                finally:
                    textpage.close()
                    page.close()
                if text:
                    docs.append(
                        Document(
                            page_content=text,
                            metadata={"source": filepath, "page": i + 1},
                        )
                    )
        finally:
            pdf.close()

        if not docs:
            raise ValueError(
                "No extractable text found in the PDF. It may be a scanned "
                "or image-only document with no selectable text."
            )

    elif filename.lower().endswith(".docx"):
        loader = Docx2txtLoader(filepath)
        docs = loader.load()

    elif filename.lower().endswith(".txt"):
        loader = TextLoader(filepath, encoding="utf-8")
        docs = loader.load()

    elif filename.lower().endswith(".pptx"):
        prs = Presentation(filepath)
        docs = []
        for i, slide in enumerate(prs.slides):
            slide_lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_lines.append(text)
            if slide_lines:
                docs.append(
                    Document(
                        page_content="\n".join(slide_lines),
                        metadata={"source": filepath, "page": i + 1},
                    )
                )
        if not docs:
            raise ValueError(
                "No readable text found in the PowerPoint file. "
                "Ensure the slides contain text content."
            )

    else:
        raise ValueError(f"Unsupported file type: {filename}")

    if not docs:
        raise ValueError(
            "No content could be extracted from the file. "
            "The file may be empty or corrupted."
        )

    # Tag every chunk with the original filename AND the owning session
    for doc in docs:
        doc.metadata["source_file"] = filename
        doc.metadata["session_id"] = session_id

    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    splits = splitter.split_documents(docs)
    if not splits:
        raise ValueError("No text chunks could be extracted from the file.")

    logger.info("'%s' → %d chunks; adding to ChromaDB (session=%s)…", filename, len(splits), session_id)
    get_db().add_documents(splits)
    logger.info("'%s' ingested successfully.", filename)


# ── RAG query ────────────────────────────────────────────────────────────────
def query_rag(query_text: str, session_id: str) -> dict:
    """Retrieve relevant chunks (scoped to this session only) and answer via the LLM."""
    logger.info("query_rag: '%s' (session=%s)", query_text[:80], session_id)

    db = get_db()
    results = db.similarity_search_with_score(
        query_text, k=3, filter={"session_id": session_id}
    )

    valid_results = [(doc, score) for doc, score in results if score <= SIMILARITY_THRESHOLD]
    logger.info(
        "similarity_search: %d total, %d below threshold (%.2f)",
        len(results), len(valid_results), SIMILARITY_THRESHOLD,
    )

    if not valid_results:
        return {
            "answer": "No data regarding this query is available in the uploaded documents.",
            "sources": [],
        }

    context_text = "\n\n---\n\n".join([doc.page_content for doc, _ in valid_results])

    prompt_template = """You are an AI assistant answering questions based on the provided Context.
Use only the information in the Context. If the Context does not contain the answer, respond with:
'I don't have enough information to answer that question.'

Context:
{context}

Question: {question}
"""
    prompt = ChatPromptTemplate.from_template(prompt_template)
    model = ChatGroq(model_name="llama-3.1-8b-instant", temperature=0)
    chain = prompt | model

    try:
        response = chain.invoke({"context": context_text, "question": query_text})
        answer_text = getattr(response, "content", None) or str(response).strip()
    except Exception as exc:
        logger.error("LLM chain error: %s", exc)
        answer_text = None

    if not answer_text:
        answer_text = "No data regarding this query is available in the uploaded documents."

    sources = []
    seen = set()
    for doc, _ in valid_results:
        snippet = doc.page_content[:200] + "..."
        if snippet not in seen:
            sources.append({
                "source": doc.metadata.get("source_file", "Unknown"),
                "page": doc.metadata.get("page", 0),
                "content": snippet,
            })
            seen.add(snippet)

    logger.info("query_rag completed. answer_len=%d sources=%d", len(answer_text), len(sources))
    return {"answer": answer_text, "sources": sources}


# ── Document management (all scoped to a session) ────────────────────────────
def list_documents(session_id: str) -> list[str]:
    """Return unique source filenames stored in ChromaDB for this session only."""
    try:
        collection = get_db()._collection
        results = collection.get(where={"session_id": session_id}, include=["metadatas"])
        if not results or not results.get("metadatas"):
            return []
        unique_files = {
            meta["source_file"]
            for meta in results["metadatas"]
            if "source_file" in meta
        }
        return sorted(unique_files)
    except Exception as exc:
        logger.error("list_documents error: %s", exc)
        return []


def delete_document(filename: str, session_id: str) -> bool:
    """Delete all chunks for the given filename, within this session only."""
    try:
        collection = get_db()._collection
        results = collection.get(
            where={"$and": [{"source_file": filename}, {"session_id": session_id}]}
        )
        if results and results.get("ids"):
            collection.delete(ids=results["ids"])
            logger.info("Deleted %d chunks for '%s' (session=%s).", len(results["ids"]), filename, session_id)
            return True
        return False
    except Exception as exc:
        logger.error("delete_document error for '%s': %s", filename, exc)
        return False


def clear_session(session_id: str) -> int:
    """Delete every chunk belonging to a session. Called on sign-out / tab close.

    Returns the number of chunks removed (0 if nothing to do — always safe to call).
    """
    try:
        collection = get_db()._collection
        results = collection.get(where={"session_id": session_id}, include=[])
        ids = results.get("ids", []) if results else []
        if ids:
            collection.delete(ids=ids)
            logger.info("clear_session: removed %d chunks for session=%s", len(ids), session_id)
        return len(ids)
    except Exception as exc:
        logger.error("clear_session error for session=%s: %s", session_id, exc)
        return 0
